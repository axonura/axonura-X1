# Copyright 2026 First Person
#
# Permission is hereby granted, free of charge, to any person obtaining a
# copy of this software and associated documentation files (the "Software"),
# to deal in the Software without restriction, including without limitation the
# rights to use, copy, modify, merge, publish, distribute, sublicense, and/or
# sell copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in
# all copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN
# THE SOFTWARE.

import csv
import math
import secrets
import tempfile
from pathlib import Path

import cv2
import filetype
import librosa
import torch
from moviepy import VideoFileClip
from PIL import Image
from transformers import (
    AutoImageProcessor,
    AutoModel,
    AutoProcessor,
)

_vidProcessor = None
_vidEncoder = None
_audProcessor = None

AUDIO_SAMPLE_RATE = 16000
MAX_VIDEO_FRAMES = 32


def _getVidProcessor():
    global _vidProcessor
    if _vidProcessor is None:
        _vidProcessor = AutoImageProcessor.from_pretrained("google/vit-base-patch16-224")
    return _vidProcessor


def _getVidEncoder():
    global _vidEncoder
    if _vidEncoder is None:
        _vidEncoder = AutoModel.from_pretrained("google/vit-base-patch16-224")
    return _vidEncoder


def _getAudProcessor():
    global _audProcessor
    if _audProcessor is None:
        _audProcessor = AutoProcessor.from_pretrained("openai/whisper-base")
    return _audProcessor


def text_gen(dataset_shard):
    for x in dataset_shard:
        text = x.get("text", "")
        if text and text.strip():
            yield text


def encode_batch(texts, tokenizer, max_len):
    texts = [t if t.strip() else " " for t in texts]
    enc = tokenizer(
        texts,
        padding="max_length",
        truncation=True,
        max_length=max_len,
        return_tensors="pt",
    )
    return enc["input_ids"]


def collate_fn(batch, tokenizer, max_len):
    input_ids = encode_batch(batch, tokenizer, max_len)
    return input_ids[:, :-1], input_ids[:, 1:]


def pad_feature_sequence(features, padValue=0.0):
    present = [f for f in features if f is not None]
    if not present:
        return None, None

    batchSize = len(features)
    timeMax = max(f.shape[0] for f in present)
    featureDim = present[0].shape[-1]
    dtype = present[0].dtype
    device = present[0].device

    padded = torch.full((batchSize, timeMax, featureDim), padValue, dtype=dtype, device=device)
    mask = torch.zeros(batchSize, timeMax, dtype=torch.bool, device=device)
    for i, feature in enumerate(features):
        if feature is None:
            continue
        length = feature.shape[0]
        padded[i, :length] = feature
        mask[i, :length] = True
    return padded, mask


def pad_audio_sequences(audios, padValue=0.0):
    present = [a for a in audios if a is not None]
    if not present:
        return None, None

    batchSize = len(audios)
    chunkMax = max(a.shape[0] for a in present)
    melBins = present[0].shape[1]
    timeSteps = present[0].shape[2]
    dtype = present[0].dtype
    device = present[0].device

    padded = torch.full(
        (batchSize, chunkMax, melBins, timeSteps), padValue, dtype=dtype, device=device
    )
    mask = torch.zeros(batchSize, chunkMax, dtype=torch.bool, device=device)
    for i, audio in enumerate(audios):
        if audio is None:
            continue
        chunks = audio.shape[0]
        padded[i, :chunks] = audio
        mask[i, :chunks] = True
    return padded, mask


def collate_multimodal(batch, tokenizer, max_len, media_token_ids):
    texts = [x.get("text") or " " for x in batch]
    enc = tokenizer(
        texts,
        padding="max_length",
        truncation=True,
        max_length=max_len,
        return_tensors="pt",
    )
    full = enc["input_ids"]

    vision_features, vision_mask = pad_feature_sequence([x.get("vision") for x in batch])
    audio_features, audio_mask = pad_audio_sequences([x.get("audio") for x in batch])

    input_ids = full[:, :-1]
    labels = full[:, 1:]
    loss_mask = torch.ones_like(labels, dtype=torch.bool)
    if tokenizer.pad_token_id is not None:
        loss_mask &= labels != tokenizer.pad_token_id
    for mediaId in media_token_ids:
        if mediaId is not None:
            loss_mask &= labels != mediaId

    return {
        "input_ids": input_ids,
        "labels": labels,
        "loss_mask": loss_mask,
        "vision_features": vision_features,
        "vision_mask": vision_mask,
        "audio_features": audio_features,
        "audio_mask": audio_mask,
    }


def batch_iterator(dataset):
    for i in range(0, len(dataset), 1000):
        yield dataset[i : i + 1000]["text"]


def top_k_sampling(logits, k=64):
    top_k_vals, top_k_indices = torch.topk(logits, k=k, dim=-1)
    top_k_logits = torch.full_like(logits, -1e9)
    top_k_logits.scatter_(-1, top_k_indices, top_k_vals)
    probs = torch.softmax(top_k_logits, dim=-1)
    return torch.multinomial(probs, num_samples=1)


def tokenizeFile(input, samplingRate=AUDIO_SAMPLE_RATE):
    input = str(input)
    kind = filetype.guess(input)
    mime = kind.mime if kind is not None else ""
    ext = Path(input).suffix.lower()

    if mime.startswith("image/"):
        return ["image", tokenizeImage(input)]
    if mime.startswith("video/"):
        frames = _tokenizeVideoFrames(input)
        audio = _extractVideoAudio(input, samplingRate)
        return ["video", _encodeFrames(frames), audio]
    if mime.startswith("audio/"):
        return ["audio", tokenizeAudio(input, samplingRate)]
    if mime == "application/pdf":
        return ["text", _extractPDF(input)]
    if "wordprocessingml" in mime or ext == ".docx":
        return ["text", _extractDocx(input)]
    if "spreadsheetml" in mime or ext == ".xlsx":
        return ["text", _extractXlsx(input)]
    if "presentationml" in mime or ext == ".pptx":
        return ["text", _extractPptx(input)]
    if mime == "application/msword" or ext == ".doc":
        raise ValueError(f"Legacy Binary Word Documents Not Supported: {input}")
    if mime == "application/vnd.ms-excel" or ext == ".xls":
        raise ValueError(f"Legacy Binary Excel Documents Not Supported: {input}")
    if mime == "application/vnd.ms-powerpoint" or ext == ".ppt":
        raise ValueError(f"Legacy Binary PowerPoint Documents Not Supported: {input}")
    if ext == ".csv":
        return ["text", _extractCsv(input)]
    if mime.startswith("text/") or ext in (".txt", ".md", ".json", ".log", ".xml", ".py", ".js", ".html"):
        return ["text", _extractPlainText(input)]
    if mime == "" and _looksLikeText(input):
        return ["text", _extractPlainText(input)]

    raise ValueError(f"Unsupported File Type Of {input} File.")


def tokenizeImage(input):
    processor = _getVidProcessor()
    encoder = _getVidEncoder()

    with Image.open(input) as image:
        if image.mode != "RGB":
            image = image.convert("RGB")
        inputs = processor(images=image, return_tensors="pt")

    with torch.no_grad():
        outputs = encoder(**inputs)

    return outputs.last_hidden_state[0]


def tokenizeAudio(input, samplingRate=AUDIO_SAMPLE_RATE):
    processor = _getAudProcessor()
    audio_array, _ = librosa.load(input, sr=samplingRate, mono=True)
    chunk_samples = 30 * samplingRate

    if len(audio_array) <= chunk_samples:
        inputs = processor(audio_array, sampling_rate=samplingRate, return_tensors="pt")
        return inputs["input_features"]

    chunks = [
        audio_array[i : i + chunk_samples]
        for i in range(0, len(audio_array), chunk_samples)
    ]
    features = [
        processor(chunk, sampling_rate=samplingRate, return_tensors="pt")["input_features"]
        for chunk in chunks
    ]
    return torch.cat(features, dim=0)


def _tokenizeVideoFrames(input, maxFrames=MAX_VIDEO_FRAMES):
    processor = _getVidProcessor()
    cap = cv2.VideoCapture(input)
    if not cap.isOpened():
        raise ValueError(f"Unable To Open Video File: {input}")

    totalFrames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
    step = max(1, int(round(fps)))
    if totalFrames > 0 and step * maxFrames < totalFrames:
        step = max(1, int(math.ceil(totalFrames / maxFrames)))

    frames = []
    index = 0
    while True:
        success, frame = cap.read()
        if not success:
            break
        if index % step == 0:
            frameRgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            pilImage = Image.fromarray(frameRgb)
            pixelValues = processor(images=pilImage, return_tensors="pt")["pixel_values"]
            frames.append(pixelValues)
            if len(frames) >= maxFrames:
                break
        index += 1
    cap.release()

    if not frames:
        raise ValueError(f"Unable To Extract Frames From Video: {input}")
    return torch.cat(frames, dim=0)


def _encodeFrames(frames):
    encoder = _getVidEncoder()
    with torch.no_grad():
        outputs = encoder(pixel_values=frames)
    return outputs.last_hidden_state


def _extractVideoAudio(input, samplingRate):
    video = None
    output = None
    try:
        video = VideoFileClip(input)
        if video.audio is None:
            return None
        output = str(Path(tempfile.gettempdir()) / f"{secrets.token_hex(16)}.wav")
        video.audio.write_audiofile(
            output,
            fps=samplingRate,
            nbytes=2,
            codec="pcm_s16le",
            ffmpeg_params=["-ac", "1"],
            logger=None,
        )
        return tokenizeAudio(output, samplingRate)
    except Exception as e:
        raise ValueError(f"Unable To Extract Audio From Video {input}: {e}") from e
    finally:
        if video is not None:
            video.close()
        if output is not None:
            Path(output).unlink(missing_ok=True)


def _extractPlainText(input):
    with open(input, "r", encoding="utf-8", errors="replace") as f:
        return f.read().strip()


def _extractCsv(input):
    parts = []
    with open(input, "r", encoding="utf-8", errors="replace", newline="") as f:
        for row in csv.reader(f):
            cells = [cell.strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extractPDF(input):
    from pypdf import PdfReader

    reader = PdfReader(input)
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    return "\n".join(parts).strip()


def _extractDocx(input):
    import docx

    document = docx.Document(input)
    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extractPptx(input):
    from pptx import Presentation

    presentation = Presentation(input)
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        parts.append(text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extractXlsx(input):
    from openpyxl import load_workbook

    workbook = load_workbook(input, read_only=True, data_only=True)
    parts = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell is not None]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        workbook.close()
    return "\n".join(parts).strip()


def _looksLikeText(input, sampleBytes=8192):
    try:
        with open(input, "rb") as f:
            sample = f.read(sampleBytes)
    except OSError:
        return False
    if not sample:
        return False
    textChars = 0
    for byte in sample:
        if byte in (9, 10, 13) or 32 <= byte <= 126 or byte >= 128:
            textChars += 1
    return textChars / len(sample) > 0.95
